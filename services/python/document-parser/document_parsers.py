"""
Document Parser Implementations
Supports PDF, Word, Excel, PowerPoint, Notion, and Confluence
"""
import os
import io
import asyncio
from typing import Dict, List, Optional, Any, Union
from pathlib import Path
import structlog
from datetime import datetime
import chardet
import aiofiles

# PDF processing
import PyPDF2

# Microsoft Office documents
from docx import Document as WordDocument
from openpyxl import load_workbook
from pptx import Presentation

# Web scraping and HTML
from bs4 import BeautifulSoup
import markdown

# Notion API
from notion_client import Client as NotionClient

# Confluence API
from atlassian import Confluence

# OCR support
try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

from doc_parser_config import config

# Import models with fallback
try:
    from models import DocumentContent, DocumentType, DocumentMetadata, ParsingError
except ImportError:
    # Create fallback models if shared models not available
    from pydantic import BaseModel
    from enum import Enum
    from typing import Optional, List, Dict, Any
    from datetime import datetime
    
    class DocumentType(str, Enum):
        PDF = "pdf"
        WORD = "word"
        EXCEL = "excel"
        POWERPOINT = "powerpoint"
        TEXT = "text"
        NOTION = "notion"
        CONFLUENCE = "confluence"
    
    class DocumentMetadata(BaseModel):
        file_name: str
        file_size: int
        document_type: DocumentType
        page_count: Optional[int] = None
        creation_date: datetime
        custom_metadata: Dict[str, Any] = {}
    
    class DocumentContent(BaseModel):
        text: str
        metadata: Optional[DocumentMetadata] = None
        sections: List[Dict[str, Any]] = []
        tables: List[Dict[str, Any]] = []
        images: List[Dict[str, Any]] = []
        links: List[str] = []
    
    class ParsingError(Exception):
        pass

logger = structlog.get_logger()


class BaseDocumentParser:
    """Base class for document parsers"""
    
    def __init__(self):
        self.supported_extensions = []
        self.max_file_size = config.max_file_size
    
    async def parse(self, file_path: str, **kwargs) -> DocumentContent:
        """Parse document and return structured content"""
        raise NotImplementedError
    
    def _validate_file(self, file_path: str) -> bool:
        """Validate file size and extension"""
        if not os.path.exists(file_path):
            raise ParsingError(f"File not found: {file_path}")
        
        file_size = os.path.getsize(file_path)
        if file_size > self.max_file_size:
            raise ParsingError(f"File too large: {file_size} bytes (max: {self.max_file_size})")
        
        file_ext = Path(file_path).suffix.lower()
        if file_ext not in self.supported_extensions:
            raise ParsingError(f"Unsupported file extension: {file_ext}")
        
        return True
    
    def _detect_encoding(self, file_path: str) -> str:
        """Detect file encoding"""
        try:
            with open(file_path, 'rb') as f:
                raw_data = f.read(10000)  # Read first 10KB
                result = chardet.detect(raw_data)
                return result.get('encoding', 'utf-8')
        except Exception:
            return 'utf-8'


class PDFParser(BaseDocumentParser):
    """PDF document parser"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.pdf']
    
    async def parse(self, file_path: str, extract_images: bool = False, **kwargs) -> DocumentContent:
        """Parse PDF document"""
        self._validate_file(file_path)
        
        logger.info("Parsing PDF document", file=file_path, extract_images=extract_images)
        
        try:
            text_content = []
            metadata = {}
            images = []
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                # Extract metadata
                if pdf_reader.metadata:
                    metadata = {
                        'title': pdf_reader.metadata.get('/Title', ''),
                        'author': pdf_reader.metadata.get('/Author', ''),
                        'subject': pdf_reader.metadata.get('/Subject', ''),
                        'creator': pdf_reader.metadata.get('/Creator', ''),
                        'producer': pdf_reader.metadata.get('/Producer', ''),
                        'creation_date': pdf_reader.metadata.get('/CreationDate', ''),
                        'modification_date': pdf_reader.metadata.get('/ModDate', '')
                    }
                
                # Extract text from each page
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            text_content.append({
                                'page': page_num + 1,
                                'content': page_text.strip()
                            })
                        
                        # Extract images if requested (basic implementation)
                        if extract_images and OCR_AVAILABLE:
                            # This would require more advanced PDF image extraction
                            pass
                            
                    except Exception as e:
                        logger.warning("Failed to extract text from page", 
                                     page=page_num + 1, error=str(e))
            
            # Combine all text
            full_text = '\n\n'.join([item['content'] for item in text_content])
            
            return DocumentContent(
                text=full_text,
                metadata=DocumentMetadata(
                    file_name=os.path.basename(file_path),
                    file_size=os.path.getsize(file_path),
                    document_type=DocumentType.PDF,
                    page_count=len(text_content),
                    creation_date=datetime.now(),
                    custom_metadata=metadata
                ),
                sections=[{
                    'title': f'Page {item["page"]}',
                    'content': item['content'],
                    'metadata': {'page_number': item['page']}
                } for item in text_content],
                images=images
            )
            
        except Exception as e:
            logger.error("PDF parsing failed", file=file_path, error=str(e))
            raise ParsingError(f"Failed to parse PDF: {str(e)}")


class WordParser(BaseDocumentParser):
    """Microsoft Word document parser"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.docx', '.doc']
    
    async def parse(self, file_path: str, **kwargs) -> DocumentContent:
        """Parse Word document"""
        self._validate_file(file_path)
        
        logger.info("Parsing Word document", file=file_path)
        
        try:
            doc = WordDocument(file_path)
            
            # Extract text and structure
            sections = []
            full_text_parts = []
            
            for i, paragraph in enumerate(doc.paragraphs):
                if paragraph.text.strip():
                    # Detect headings based on style
                    is_heading = paragraph.style.name.startswith('Heading')
                    
                    sections.append({
                        'title': f'Paragraph {i+1}' if not is_heading else paragraph.text.strip(),
                        'content': paragraph.text.strip(),
                        'metadata': {
                            'paragraph_number': i + 1,
                            'style': paragraph.style.name,
                            'is_heading': is_heading
                        }
                    })
                    full_text_parts.append(paragraph.text.strip())
            
            # Extract tables
            tables = []
            for table_num, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                
                tables.append({
                    'table_number': table_num + 1,
                    'data': table_data
                })
            
            # Get document properties
            props = doc.core_properties
            metadata = {
                'title': props.title or '',
                'author': props.author or '',
                'subject': props.subject or '',
                'keywords': props.keywords or '',
                'created': props.created.isoformat() if props.created else '',
                'modified': props.modified.isoformat() if props.modified else ''
            }
            
            return DocumentContent(
                text='\n\n'.join(full_text_parts),
                metadata=DocumentMetadata(
                    file_name=os.path.basename(file_path),
                    file_size=os.path.getsize(file_path),
                    document_type=DocumentType.WORD,
                    creation_date=datetime.now(),
                    custom_metadata=metadata
                ),
                sections=sections,
                tables=tables
            )
            
        except Exception as e:
            logger.error("Word parsing failed", file=file_path, error=str(e))
            raise ParsingError(f"Failed to parse Word document: {str(e)}")


class ExcelParser(BaseDocumentParser):
    """Microsoft Excel parser"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.xlsx', '.xls']
    
    async def parse(self, file_path: str, **kwargs) -> DocumentContent:
        """Parse Excel document"""
        self._validate_file(file_path)
        
        logger.info("Parsing Excel document", file=file_path)
        
        try:
            workbook = load_workbook(file_path, data_only=True)
            
            sections = []
            tables = []
            full_text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # Extract data from sheet
                sheet_data = []
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        row_data = [str(cell) if cell is not None else '' for cell in row]
                        sheet_data.append(row_data)
                
                if sheet_data:
                    tables.append({
                        'sheet_name': sheet_name,
                        'data': sheet_data
                    })
                    
                    # Convert to text representation
                    sheet_text = f"Sheet: {sheet_name}\n"
                    for row in sheet_data[:100]:  # Limit to first 100 rows
                        sheet_text += '\t'.join(row) + '\n'
                    
                    sections.append({
                        'title': f'Sheet: {sheet_name}',
                        'content': sheet_text,
                        'metadata': {
                            'sheet_name': sheet_name,
                            'row_count': len(sheet_data)
                        }
                    })
                    full_text_parts.append(sheet_text)
            
            workbook.close()
            
            return DocumentContent(
                text='\n\n'.join(full_text_parts),
                metadata=DocumentMetadata(
                    file_name=os.path.basename(file_path),
                    file_size=os.path.getsize(file_path),
                    document_type=DocumentType.EXCEL,
                    creation_date=datetime.now(),
                    custom_metadata={'sheet_count': len(workbook.sheetnames)}
                ),
                sections=sections,
                tables=tables
            )
            
        except Exception as e:
            logger.error("Excel parsing failed", file=file_path, error=str(e))
            raise ParsingError(f"Failed to parse Excel document: {str(e)}")


class PowerPointParser(BaseDocumentParser):
    """Microsoft PowerPoint parser"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.pptx', '.ppt']
    
    async def parse(self, file_path: str, **kwargs) -> DocumentContent:
        """Parse PowerPoint document"""
        self._validate_file(file_path)
        
        logger.info("Parsing PowerPoint document", file=file_path)
        
        try:
            prs = Presentation(file_path)
            
            sections = []
            full_text_parts = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text_parts = []
                slide_title = f"Slide {slide_num + 1}"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text_parts.append(shape.text.strip())
                        
                        # Use first text as slide title if it looks like a title
                        if not slide_title.startswith("Slide") and len(shape.text.strip()) < 100:
                            slide_title = shape.text.strip()[:50] + "..."
                
                if slide_text_parts:
                    slide_content = '\n'.join(slide_text_parts)
                    sections.append({
                        'title': slide_title,
                        'content': slide_content,
                        'metadata': {
                            'slide_number': slide_num + 1,
                            'shape_count': len(slide.shapes)
                        }
                    })
                    full_text_parts.append(f"{slide_title}\n{slide_content}")
            
            return DocumentContent(
                text='\n\n'.join(full_text_parts),
                metadata=DocumentMetadata(
                    file_name=os.path.basename(file_path),
                    file_size=os.path.getsize(file_path),
                    document_type=DocumentType.POWERPOINT,
                    creation_date=datetime.now(),
                    custom_metadata={'slide_count': len(prs.slides)}
                ),
                sections=sections
            )
            
        except Exception as e:
            logger.error("PowerPoint parsing failed", file=file_path, error=str(e))
            raise ParsingError(f"Failed to parse PowerPoint document: {str(e)}")


class NotionParser:
    """Notion document parser"""
    
    def __init__(self, notion_token: Optional[str] = None):
        self.notion_token = notion_token or config.notion_token
        if self.notion_token:
            self.client = NotionClient(auth=self.notion_token)
        else:
            self.client = None
            logger.warning("Notion token not provided - Notion parsing disabled")
    
    async def parse_page(self, page_id: str, **kwargs) -> DocumentContent:
        """Parse Notion page"""
        if not self.client:
            raise ParsingError("Notion client not initialized - check token")
        
        logger.info("Parsing Notion page", page_id=page_id)
        
        try:
            # Get page info
            page = await asyncio.to_thread(self.client.pages.retrieve, page_id)
            
            # Get page blocks
            blocks = await asyncio.to_thread(
                self.client.blocks.children.list, 
                block_id=page_id
            )
            
            sections = []
            full_text_parts = []
            
            def extract_block_text(block):
                """Extract text from a block"""
                block_type = block.get('type', '')
                
                if block_type in ['paragraph', 'heading_1', 'heading_2', 'heading_3']:
                    rich_text = block.get(block_type, {}).get('rich_text', [])
                    return ''.join([text.get('plain_text', '') for text in rich_text])
                elif block_type == 'bulleted_list_item':
                    rich_text = block.get('bulleted_list_item', {}).get('rich_text', [])
                    return '• ' + ''.join([text.get('plain_text', '') for text in rich_text])
                elif block_type == 'numbered_list_item':
                    rich_text = block.get('numbered_list_item', {}).get('rich_text', [])
                    return '1. ' + ''.join([text.get('plain_text', '') for text in rich_text])
                elif block_type == 'code':
                    rich_text = block.get('code', {}).get('rich_text', [])
                    return '```\n' + ''.join([text.get('plain_text', '') for text in rich_text]) + '\n```'
                elif block_type == 'quote':
                    rich_text = block.get('quote', {}).get('rich_text', [])
                    return '> ' + ''.join([text.get('plain_text', '') for text in rich_text])
                
                return ''
            
            for block in blocks.get('results', []):
                text = extract_block_text(block)
                if text.strip():
                    block_type = block.get('type', 'unknown')
                    sections.append({
                        'title': f"{block_type.replace('_', ' ').title()}",
                        'content': text.strip(),
                        'metadata': {
                            'block_type': block_type,
                            'block_id': block.get('id', '')
                        }
                    })
                    full_text_parts.append(text.strip())
            
            # Extract page title
            title = page.get('properties', {}).get('title', {})
            page_title = ''
            if title.get('title'):
                page_title = ''.join([t.get('plain_text', '') for t in title['title']])
            
            return DocumentContent(
                text='\n\n'.join(full_text_parts),
                metadata=DocumentMetadata(
                    file_name=f"notion_page_{page_id}.md",
                    file_size=len('\n\n'.join(full_text_parts).encode('utf-8')),
                    document_type=DocumentType.NOTION,
                    creation_date=datetime.now(),
                    custom_metadata={
                        'page_id': page_id,
                        'title': page_title,
                        'url': page.get('url', ''),
                        'created_time': page.get('created_time', ''),
                        'last_edited_time': page.get('last_edited_time', '')
                    }
                ),
                sections=sections
            )
            
        except Exception as e:
            logger.error("Notion parsing failed", page_id=page_id, error=str(e))
            raise ParsingError(f"Failed to parse Notion page: {str(e)}")


class ConfluenceParser:
    """Confluence document parser"""
    
    def __init__(self, confluence_url: Optional[str] = None, confluence_token: Optional[str] = None):
        self.confluence_url = confluence_url or config.confluence_url
        self.confluence_token = confluence_token or config.confluence_token
        
        if self.confluence_url and self.confluence_token:
            self.client = Confluence(
                url=self.confluence_url,
                token=self.confluence_token
            )
        else:
            self.client = None
            logger.warning("Confluence credentials not provided - Confluence parsing disabled")
    
    async def parse_page(self, page_id: str, **kwargs) -> DocumentContent:
        """Parse Confluence page"""
        if not self.client:
            raise ParsingError("Confluence client not initialized - check credentials")
        
        logger.info("Parsing Confluence page", page_id=page_id)
        
        try:
            # Get page content
            page = await asyncio.to_thread(
                self.client.get_page_by_id,
                page_id,
                expand='body.storage,space,version,ancestors'
            )
            
            # Extract HTML content
            html_content = page.get('body', {}).get('storage', {}).get('value', '')
            
            # Parse HTML with BeautifulSoup
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Extract structured content
            sections = []
            full_text_parts = []
            
            # Find headings and content
            for element in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'ul', 'ol', 'code', 'pre']):
                if element.get_text().strip():
                    element_type = element.name
                    content = element.get_text().strip()
                    
                    if element_type.startswith('h'):
                        title = content
                    else:
                        title = f"{element_type.upper()} Content"
                    
                    sections.append({
                        'title': title,
                        'content': content,
                        'metadata': {
                            'element_type': element_type,
                            'tag': element.name
                        }
                    })
                    full_text_parts.append(content)
            
            return DocumentContent(
                text='\n\n'.join(full_text_parts),
                metadata=DocumentMetadata(
                    file_name=f"confluence_page_{page_id}.html",
                    file_size=len(html_content.encode('utf-8')),
                    document_type=DocumentType.CONFLUENCE,
                    creation_date=datetime.now(),
                    custom_metadata={
                        'page_id': page_id,
                        'title': page.get('title', ''),
                        'space': page.get('space', {}).get('name', ''),
                        'version': page.get('version', {}).get('number', 0),
                        'url': page.get('_links', {}).get('webui', '')
                    }
                ),
                sections=sections
            )
            
        except Exception as e:
            logger.error("Confluence parsing failed", page_id=page_id, error=str(e))
            raise ParsingError(f"Failed to parse Confluence page: {str(e)}")


class TextParser(BaseDocumentParser):
    """Plain text and markdown parser"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.txt', '.md', '.rst']
    
    async def parse(self, file_path: str, **kwargs) -> DocumentContent:
        """Parse text document"""
        self._validate_file(file_path)
        
        logger.info("Parsing text document", file=file_path)
        
        try:
            encoding = self._detect_encoding(file_path)
            
            async with aiofiles.open(file_path, 'r', encoding=encoding) as f:
                content = await f.read()
            
            file_ext = Path(file_path).suffix.lower()
            
            # Process markdown
            if file_ext == '.md':
                html_content = markdown.markdown(content)
                soup = BeautifulSoup(html_content, 'html.parser')
                
                sections = []
                for heading in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
                    # Find content after heading until next heading
                    section_content = []
                    for sibling in heading.next_siblings:
                        if sibling.name and sibling.name.startswith('h'):
                            break
                        if sibling.get_text().strip():
                            section_content.append(sibling.get_text().strip())
                    
                    sections.append({
                        'title': heading.get_text().strip(),
                        'content': '\n'.join(section_content),
                        'metadata': {
                            'heading_level': int(heading.name[1]),
                            'type': 'markdown_section'
                        }
                    })
            else:
                # Plain text - split by paragraphs
                paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
                sections = [{
                    'title': f'Paragraph {i+1}',
                    'content': paragraph,
                    'metadata': {'paragraph_number': i+1}
                } for i, paragraph in enumerate(paragraphs)]
            
            return DocumentContent(
                text=content,
                metadata=DocumentMetadata(
                    file_name=os.path.basename(file_path),
                    file_size=os.path.getsize(file_path),
                    document_type=DocumentType.TEXT,
                    creation_date=datetime.now(),
                    custom_metadata={'encoding': encoding}
                ),
                sections=sections
            )
            
        except Exception as e:
            logger.error("Text parsing failed", file=file_path, error=str(e))
            raise ParsingError(f"Failed to parse text document: {str(e)}")